#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT图片插入工具 - 图形界面版本 更新记录
v4.1.1新功能：文本框形式插入文本、一键插字、字体、字间距调整，所有文本从日志文件统一读取
v4.1.2新功能：文本不再从单一日志文件读取（依据大数字行号），改为关键词检索文件
"""

import tkinter as tk
from tkinter import ttk, filedialog, messagebox, Canvas, simpledialog
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import json
from PIL import Image, ImageTk, ImageDraw, ImageFont
from copy import deepcopy
from datetime import datetime
import re
from decimal import Decimal, getcontext, ROUND_HALF_UP


def format_number(value):
    """将数字格式化为保留三位有效数字的普通小数点格式"""
    try:
        # 尝试解析为浮点数
        num = float(str(value).strip())

        # 处理特殊值
        if not isfinite(num):
            return str(value)

        # 设置decimal精度
        getcontext().prec = 10

        # 转换为Decimal
        dec_num = Decimal(value)

        # 计算数量级
        if dec_num == 0:
            return "0.00"

        # 获取科学计数法表示
        sci_notation = f"{abs(dec_num):.3e}"
        base, exponent = sci_notation.split('e')
        exponent = int(exponent)

        # 根据数量级决定格式
        if exponent >= -1 and exponent <= 2:
            # 小范围数字，保留3位有效数字并转为普通格式
            target_decimal_places = 3 - (exponent + 1)
            if target_decimal_places > 0:
                rounded = dec_num.quantize(Decimal(f"1.{'0' * min(target_decimal_places, 10)}"), rounding=ROUND_HALF_UP)
                return f"{rounded:.{target_decimal_places}f}".rstrip('0').rstrip('.')
            else:
                return f"{int(dec_num)}"
        else:
            # 大范围数字，使用科学计数法但保留3位有效数字
            decimal_places = 2
            rounded = dec_num.quantize(Decimal(f"1e{exponent - decimal_places}"), rounding=ROUND_HALF_UP)
            return f"{rounded:.{decimal_places}e}".replace('e', 'E')

    except (ValueError, ZeroDivisionError):
        # 如果不是数字，返回原值
        try:
            # 检查是否已经是类似科学计数法的形式
            if 'E' in str(value) or 'e' in str(value):
                return str(value)
            return str(value).strip()
        except:
            return str(value)


def isfinite(x):
    """检查数字是否有限"""
    try:
        import math
        return math.isfinite(x)
    except:
        return True


def format_text(text):
    """格式化文本数据（如果是数字则格式化，否则返回原值）"""
    try:
        # 尝试解析为浮点数进行格式化
        return format_number(text)
    except:
        # 不是数字则返回原值
        return str(text).strip()


# 配置文件路径
CONFIG_DIR = os.path.join(os.path.expanduser("~"), ".ppt_image_inserter")
MODES_FILE = os.path.join(CONFIG_DIR, "custom_modes.json")

# 默认模板文件路径（所有模板都放在一个PPT中，幻灯片索引对应布局模式）
DEFAULT_TEMPLATE_FILE = "templates/Templates.pptx"

# 确保配置目录存在
if not os.path.exists(CONFIG_DIR):
    os.makedirs(CONFIG_DIR)


class RoundedButton(tk.Canvas):
    """圆角按钮类"""

    def __init__(self, parent, text, command=None, bg='#E8F4E8', hover_bg='#D4E8D4',
                 text_color='#333', font=("微软雅黑", 10), width=120, height=35,
                 corner_radius=8, cursor='hand2', **kwargs):
        self.text = text
        self.command = command
        self.bg = bg
        self.hover_bg = hover_bg
        self.text_color = text_color
        self.font = font
        self.corner_radius = corner_radius

        # 移除不需要的参数
        kwargs.pop('relief', None)
        kwargs.pop('padx', None)
        kwargs.pop('pady', None)
        kwargs.pop('borderwidth', None)
        kwargs.pop('highlightthickness', None)
        kwargs.pop('highlightbackground', None)

        super().__init__(parent, width=width, height=height, bg='white',
                        highlightthickness=0, cursor=cursor, **kwargs)

        self.is_hovered = False
        self.text_id = None
        self.rect_id = None
        self.draw_button()

        # 绑定事件
        self.bind('<Enter>', self.on_enter)
        self.bind('<Leave>', self.on_leave)
        self.bind('<Button-1>', lambda e: self.on_click())

    def draw_button(self):
        """绘制圆角按钮"""
        self.delete("all")
        width = self.winfo_width() if self.winfo_width() > 1 else int(self['width'])
        height = self.winfo_height() if self.winfo_height() > 1 else int(self['height'])

        # 绘制圆角矩形
        x1, y1 = 2, 2
        x2, y2 = width - 2, height - 2

        # 创建圆角矩形路径
        self.rect_id = self.create_rounded_rect(x1, y1, x2, y2, self.corner_radius,
                                               fill=self.current_bg, outline='')

        # 绘制文字
        self.text_id = self.create_text(width / 2, height / 2, text=self.text,
                                       fill=self.text_color, font=self.font)

    def create_rounded_rect(self, x1, y1, x2, y2, r, **kwargs):
        """创建圆角矩形"""
        points = [
            x1 + r, y1,
            x2 - r, y1,
            x2, y1,
            x2, y1 + r,
            x2, y2 - r,
            x2, y2,
            x2 - r, y2,
            x1 + r, y2,
            x1, y2,
            x1, y2 - r,
            x1, y1 + r,
            x1, y1
        ]
        return self.create_polygon(points, smooth=True, **kwargs)

    @property
    def current_bg(self):
        """当前背景色"""
        return self.hover_bg if self.is_hovered else self.bg

    def on_enter(self, event):
        """鼠标进入"""
        self.is_hovered = True
        self.itemconfig(self.rect_id, fill=self.current_bg)

    def on_leave(self, event):
        """鼠标离开"""
        self.is_hovered = False
        self.itemconfig(self.rect_id, fill=self.current_bg)

    def on_click(self):
        """点击事件"""
        if self.command:
            self.command()

    def configure(self, text=None, bg=None, **kwargs):
        """更新按钮属性"""
        if text is not None:
            self.text = text
            self.itemconfig(self.text_id, text=text)
        if bg is not None:
            self.bg = bg
            self.itemconfig(self.rect_id, fill=self.current_bg)
        super().configure(**kwargs)


# 公共函数：将PPT单位转换为厘米
def convert_to_cm(ppt_unit):
    """将PPT的emu单位转换为厘米"""
    if hasattr(ppt_unit, 'cm'):
        return ppt_unit.cm
    elif hasattr(ppt_unit, 'inches'):
        return ppt_unit.inches * 2.54
    elif hasattr(ppt_unit, 'pt'):
        return ppt_unit.pt * 0.0352778
    else:
        # emu单位转换为厘米（1 emu = 914400 emu per inch）
        emu_value = float(ppt_unit)
        return (emu_value / 914400) * 2.54



def load_custom_modes():
    """加载自定义贴图模式（包含模板信息和文本布局）"""
    default_modes = {
        "自定义": {
            "description": "手动配置每张图片的位置和大小",
            "template_file": None,  # 无模板
            "slide_index": 0,
            "layouts": [],
            "text_layouts": []  # 新增：文本布局列表
        }
    }

    if os.path.exists(MODES_FILE):
        try:
            with open(MODES_FILE, 'r', encoding='utf-8') as f:
                modes = json.load(f)
                # 确保自定义模式存在
                if "自定义" not in modes:
                    modes["自定义"] = default_modes["自定义"]
                # 兼容旧版本数据
                for name, mode in modes.items():
                    # 兼容旧版本数据（没有template_file和slide_index的情况）
                    if "template_file" not in mode:
                        mode["template_file"] = None
                    if "slide_index" not in mode:
                        mode["slide_index"] = 0
                    # 兼容旧版本数据（没有text_layouts的情况）
                    if "text_layouts" not in mode:
                        mode["text_layouts"] = []
                return modes
        except:
            return default_modes
    return default_modes


def save_custom_modes(modes):
    """保存自定义贴图模式"""
    try:
        with open(MODES_FILE, 'w', encoding='utf-8') as f:
            json.dump(modes, f, ensure_ascii=False, indent=2)
        return True
    except Exception as e:
        messagebox.showerror("错误", f"保存模式失败: {str(e)}")
        return False


class ImageEntry:
    """图片条目类，用于管理单个图片的配置"""

    def __init__(self, parent, on_delete, index=0, app_master=None):
        self.frame = tk.Frame(parent, bg='white', highlightbackground='#E0E0E0', highlightthickness=1)
        self.frame.pack(fill=tk.X, padx=5, pady=5)
        self.on_delete = on_delete
        self.index = index
        self.app_master = app_master  # 引用主界面，用于获取工作路径

        # 图片文件名（不含路径）
        self.image_filename = tk.StringVar()

        # 第一行：图片文件名和删除按钮
        path_frame = tk.Frame(self.frame, bg='white')
        path_frame.pack(fill=tk.X, padx=5, pady=5)

        tk.Label(path_frame, text=f"图片{index+1}:", bg='white', font=("微软雅黑", 9)).pack(side=tk.LEFT)
        tk.Entry(path_frame, textvariable=self.image_filename, width=35,
                font=("微软雅黑", 9)).pack(side=tk.LEFT, padx=5)
        RoundedButton(path_frame, text="删除", command=self.delete_self,
                     bg='#FFE0E0', hover_bg='#FFD0D0', font=("微软雅黑", 8),
                     width=70, height=28, corner_radius=10).pack(side=tk.LEFT)

        # 位置和尺寸设置
        settings_frame = tk.Frame(self.frame, bg='white')
        settings_frame.pack(fill=tk.X, padx=5, pady=5)

        # 左边距
        tk.Label(settings_frame, text="左(cm):", bg='white', font=("微软雅黑", 8)).grid(row=0, column=0, padx=3)
        self.left_var = tk.StringVar(value="2")
        tk.Entry(settings_frame, textvariable=self.left_var, width=8,
                font=("微软雅黑", 8)).grid(row=0, column=1, padx=3)

        # 上边距
        tk.Label(settings_frame, text="上(cm):", bg='white', font=("微软雅黑", 8)).grid(row=0, column=2, padx=3)
        self.top_var = tk.StringVar(value="2")
        tk.Entry(settings_frame, textvariable=self.top_var, width=8,
                font=("微软雅黑", 8)).grid(row=0, column=3, padx=3)

        # 宽度
        tk.Label(settings_frame, text="宽(cm):", bg='white', font=("微软雅黑", 8)).grid(row=0, column=4, padx=3)
        self.width_var = tk.StringVar(value="")  # 默认为空，保持比例
        tk.Entry(settings_frame, textvariable=self.width_var, width=8,
                font=("微软雅黑", 8)).grid(row=0, column=5, padx=3)

        # 高度
        tk.Label(settings_frame, text="高(cm):", bg='white', font=("微软雅黑", 8)).grid(row=0, column=6, padx=3)
        self.height_var = tk.StringVar(value="8")  # 默认高度8cm
        tk.Entry(settings_frame, textvariable=self.height_var, width=8,
                font=("微软雅黑", 8)).grid(row=0, column=7, padx=3)

    def delete_self(self):
        """删除当前条目"""
        self.frame.destroy()
        self.on_delete(self)

    def set_layout(self, layout):
        """设置布局参数（从布局模式加载）"""
        self.left_var.set(f"{layout.get('left', 2):.2f}")
        self.top_var.set(f"{layout.get('top', 2):.2f}")
        self.image_filename.set(layout.get('filename', ''))
        # 宽度和高度可以为空
        if layout.get('width'):
            self.width_var.set(f"{layout.get('width'):.2f}")
        else:
            self.width_var.set("")
        if layout.get('height'):
            self.height_var.set(f"{layout.get('height'):.2f}")
        else:
            self.height_var.set("")

    def get_config(self):
        """获取当前配置"""
        try:
            config = {
                "left": float(self.left_var.get()),
                "top": float(self.top_var.get()),
            }
            # 文件名从image_filename获取
            config["filename"] = self.image_filename.get()

            if self.width_var.get().strip():
                config["width"] = float(self.width_var.get())
            if self.height_var.get().strip():
                config["height"] = float(self.height_var.get())

            return config
        except ValueError as e:
            raise ValueError(f"配置错误: {str(e)}")


class TextEntry:
    """文本条目类，用于管理单个文本卡片的配置"""

    def __init__(self, parent, on_delete, index=0, app_master=None):
        self.frame = tk.Frame(parent, bg='white', highlightbackground='#FFE0E0', highlightthickness=1)
        self.frame.pack(fill=tk.X, padx=5, pady=5)
        self.on_delete = on_delete
        self.index = index
        self.app_master = app_master  # 引用主界面，用于获取工作路径

        # 第一行：标题、关键词和删除按钮
        path_frame = tk.Frame(self.frame, bg='white')
        path_frame.pack(fill=tk.X, padx=5, pady=5)

        tk.Label(path_frame, text=f"文本{index+1}:", bg='white', font=("微软雅黑", 9), fg='#8B4513').pack(side=tk.LEFT)

        # 关键词输入框
        tk.Label(path_frame, text="关键词:", bg='white', font=("微软雅黑", 8)).pack(side=tk.LEFT, padx=(15, 5))
        self.keyword_var = tk.StringVar(value="")
        tk.Entry(path_frame, textvariable=self.keyword_var, width=15,
                font=("微软雅黑", 8)).pack(side=tk.LEFT, padx=(0, 10))

        RoundedButton(path_frame, text="删除", command=self.delete_self,
                     bg='#FFE0E0', hover_bg='#FFD0D0', font=("微软雅黑", 8),
                     width=70, height=28, corner_radius=10).pack(side=tk.LEFT, padx=(5, 0))

        # 第二行：文件行号、文件列、文本框位置
        settings_frame = tk.Frame(self.frame, bg='white')
        settings_frame.pack(fill=tk.X, padx=5, pady=5)

        # 文件行号
        tk.Label(settings_frame, text="行:", bg='white', font=("微软雅黑", 8)).grid(row=0, column=0, padx=3, sticky=tk.W)
        self.line_number_var = tk.StringVar(value="1")
        tk.Entry(settings_frame, textvariable=self.line_number_var, width=5,
                font=("微软雅黑", 8)).grid(row=0, column=1, padx=3)

        # 文件列（多个列用逗号分隔）
        tk.Label(settings_frame, text="列:", bg='white', font=("微软雅黑", 8)).grid(row=0, column=2, padx=3, sticky=tk.W)
        self.file_cols_var = tk.StringVar(value="1")
        tk.Entry(settings_frame, textvariable=self.file_cols_var, width=5,
                font=("微软雅黑", 8)).grid(row=0, column=3, padx=3)

        # 左坐标
        tk.Label(settings_frame, text="左(cm):", bg='white', font=("微软雅黑", 8)).grid(row=0, column=4, padx=3, sticky=tk.W)
        self.left_var = tk.StringVar(value="2")
        tk.Entry(settings_frame, textvariable=self.left_var, width=5,
                font=("微软雅黑", 8)).grid(row=0, column=5, padx=3)

        # 上坐标
        tk.Label(settings_frame, text="上(cm):", bg='white', font=("微软雅黑", 8)).grid(row=0, column=6, padx=3, sticky=tk.W)
        self.top_var = tk.StringVar(value="2")
        tk.Entry(settings_frame, textvariable=self.top_var, width=5,
                font=("微软雅黑", 8)).grid(row=0, column=7, padx=3)

    def delete_self(self):
        """删除当前条目"""
        self.frame.destroy()
        self.on_delete(self)

    def set_layout(self, layout):
        """设置布局参数（从布局模式加载）"""
        self.line_number_var.set(str(layout.get('line_number', 1)))
        self.file_cols_var.set(layout.get('file_cols', '1'))
        self.left_var.set(str(layout.get('left', 2)))
        self.top_var.set(str(layout.get('top', 2)))
        self.keyword_var.set(layout.get('keyword', ''))

    def get_config(self):
        """获取当前配置"""
        try:
            config = {
                "line_number": int(self.line_number_var.get()),
                "file_cols": self.file_cols_var.get(),  # 字符串，可能有多个列号
                "left": float(self.left_var.get()),
                "top": float(self.top_var.get()),
                "keyword": self.keyword_var.get().strip()  # 添加关键词
            }
            return config
        except ValueError as e:
            raise ValueError(f"配置错误: {str(e)}")


class LayoutPreviewCanvas(tk.Canvas):
    """布局预览画布"""

    def __init__(self, parent, width=300, height=200, use_letters=False):
        super().__init__(parent, width=width, height=height, bg='white', relief=tk.FLAT, borderwidth=0)
        self.canvas_width = width
        self.canvas_height = height
        self.ppt_width = 66.69  # PPT宽度（厘米）- 用户自定义尺寸
        self.ppt_height = 37.27  # PPT高度（厘米）- 用户自定义尺寸
        self.use_letters = use_letters  # 是否使用字母命名

    def draw_layout(self, layouts, text_layouts=None):
        """绘制布局预览（包含图片和文本框）"""
        self.delete("all")

        if not layouts and not text_layouts:
            self.create_text(self.canvas_width/2, self.canvas_height/2,
                           text="自定义模式\n手动配置",
                           font=("微软雅黑", 10), fill="gray")
            return

        # 获取实际画布尺寸，如果画布还未布局则使用初始尺寸
        actual_width = max(self.winfo_width(), self.canvas_width)
        actual_height = max(self.winfo_height(), self.canvas_height)
        margin = 10  # 边距，防止内容贴边被裁剪
        content_width = actual_width - 2 * margin
        content_height = actual_height - 2 * margin

        # 不绘制PPT边框，去掉灰色方框

        # 计算缩放比例（以较小的维为基准，保持宽高比）
        scale = min(content_width / self.ppt_width, content_height / self.ppt_height)

        # 计算居中偏移
        total_width = self.ppt_width * scale
        total_height = self.ppt_height * scale
        offset_x = (actual_width - total_width) / 2
        offset_y = (actual_height - total_height) / 2

        # 绘制每个图片位置
        for i, layout in enumerate(layouts):
            left_cm = layout.get("left", 0)
            top_cm = layout.get("top", 0)
            left = left_cm * scale + offset_x
            top = top_cm * scale + offset_y

            # 获取宽度和高度，如果没有则使用默认值
            width_cm = layout.get("width", None)
            height_cm = layout.get("height", 8)  # 默认高度8cm

            # 如果只有高度，假设宽高比为 4:3
            if width_cm is None and height_cm is not None:
                width_cm = height_cm * 4 / 3
            elif width_cm is not None and height_cm is None:
                height_cm = width_cm * 3 / 4
            elif width_cm is None and height_cm is None:
                width_cm = 10
                height_cm = 8

            width = width_cm * scale
            height = height_cm * scale

            # 绘制矩形框
            x1, y1 = left, top
            x2, y2 = left + width, top + height

            # 使用不同颜色：原有6色 + 4色低饱和度
            colors = [
                "#FFE5E5", "#E5F5FF", "#E5FFE5", "#FFF5E5",  # 原有1-4色
                "#F5E5FF", "#FFE5F5",  # 原有5-6色
                "#F0E5E5", "#E5F0F5", "#F5F0E5", "#F5E5F0"  # 新增4色（低饱和度）
            ]
            fill_color = colors[i % len(colors)]

            self.create_rectangle(x1, y1, x2, y2,
                                fill=fill_color,
                                outline="#666",
                                width=2)

            # 添加标签（自适应字体大小）
            center_x = (x1 + x2) / 2
            center_y = (y1 + y2) / 2

            # 根据矩形大小调整字体
            font_size = max(8, min(12, int(min(width, height) / 5)))

            # 根据设置显示字母或数字
            if self.use_letters:
                # 使用字母：a, b, c...
                label = chr(ord('a') + i)
            else:
                # 使用数字：图片1, 图片2, 图片3...
                label = f"图片{i+1}"

            self.create_text(center_x, center_y,
                           text=label,
                           font=("微软雅黑", font_size, "bold"),
                           fill="#333")

        # 绘制文本框（浅灰色）
        if text_layouts:
            for i, text_layout in enumerate(text_layouts):
                left_cm = text_layout.get("left", 0)
                top_cm = text_layout.get("top", 0)
                left = left_cm * scale + offset_x
                top = top_cm * scale + offset_y

                # 文本框默认宽度和高度
                width_cm = 5  # 默认宽度5cm
                height_cm = 1  # 默认高度1cm
                width = width_cm * scale
                height = height_cm * scale

                # 绘制矩形框（浅灰色）
                x1, y1 = left, top
                x2, y2 = left + width, top + height

                self.create_rectangle(x1, y1, x2, y2,
                                    fill="#E0E0E0",
                                    outline="#999",
                                    width=2)

                # 添加标签
                center_x = (x1 + x2) / 2
                center_y = (y1 + y2) / 2
                self.create_text(center_x, center_y,
                               text=f"文本{i+1}",
                               font=("微软雅黑", 9, "bold"),
                               fill="#666")


class PPTImageInserterGUI:
    """PPT图片插入工具GUI主类"""

    def __init__(self, root):
        self.root = root
        self.root.title("PPT自动化工具 v4.1.2")
        self.root.geometry("1200x600")  # 调整窗口尺寸
        self.root.minsize(1000, 600)
        self.root.configure(bg='white')

        self.image_entries = []
        self.text_entries = []
        self.preset_modes = load_custom_modes()
        self.current_mode = tk.StringVar(value="自定义")
        self.list_info_var = tk.StringVar(value="（可上下滚动）")
        self.preview_info_var = tk.StringVar(value="")
        self.info_hint = tk.StringVar(value="（提示）选择布局模式后，程序将自动加载模板和图片布局")
        self.template_path = tk.StringVar()
        self.template_filename = tk.StringVar()
        self.work_path = tk.StringVar()  # 工作路径（原output_path）

        self.create_widgets()

    def create_widgets(self):
        """创建界面组件"""

        # 设置样式
        style = ttk.Style()
        style.configure('TFrame', background='white')
        style.configure('TLabel', background='white')
        style.configure('TLabelframe', background='white')
        style.configure('TLabelframe.Label', background='white')

        # 自定义滚动条样式（更宽，浅灰色，圆润）
        style.theme_use('default')  # 改用默认主题，clam在Windows上可能有问题
        style.configure('Vertical.TScrollbar',
                        background='#D0D0D0',
                        troughcolor='#F0F0F0',
                        bordercolor='#F0F0F0',
                        arrowsize=0,
                        gripcount=0,
                        relief=tk.RAISED,
                        thickness=30)  # 使用thickness参数
        style.map('Vertical.TScrollbar',
                 background=[('active', '#BDBDBD'), ('pressed', '#9E9E9E')])
        style.configure('Horizontal.TScrollbar',
                        background='#E8E8E8',
                        troughcolor='#F5F5F5',
                        bordercolor='#F5F5F5',
                        arrowsize=0)

        # 主容器，白色背景
        main_container = tk.Frame(self.root, bg='white')
        main_container.pack(fill=tk.BOTH, expand=True, padx=20, pady=15)

        # ========== 顶部区域：布局模式选择（核心功能） ==========
        top_section = tk.Frame(main_container, bg='white')
        top_section.pack(fill=tk.X, pady=(0, 15))

        # 第一行：布局模式选择 + 模板文件
        row1 = tk.Frame(top_section, bg='white')
        row1.pack(fill=tk.X, pady=(0, 15))

        # 布局模式选择
        tk.Label(row1, text="布局模式", font=("微软雅黑", 11, "bold"),
                bg='white').pack(side=tk.LEFT, padx=(0, 10))

        self.mode_combo = ttk.Combobox(row1, textvariable=self.current_mode,
                                       values=list(self.preset_modes.keys()),
                                       state="readonly", width=18, font=("微软雅黑", 10))
        self.mode_combo.pack(side=tk.LEFT, padx=(0, 10))
        self.mode_combo.bind("<<ComboboxSelected>>", self.on_mode_change)

        # 合并后的模板按钮
        RoundedButton(row1, text="设置模板文件", command=self.browse_template,
                     bg='#FFF3E0', hover_bg='#FFE0C0', font=("微软雅黑", 9, "bold"),
                     width=120, height=32, corner_radius=10).pack(side=tk.LEFT, padx=(0, 10))

        RoundedButton(row1, text="删除布局模式", command=self.delete_layout_mode,
                     bg='#FFE0E0', hover_bg='#FFD0D0', font=("微软雅黑", 9, "bold"),
                     width=120, height=32, corner_radius=10).pack(side=tk.LEFT, padx=(0, 10))

        # 模板文件名显示（只显示文件名，不显示完整路径）
        self.template_filename = tk.StringVar()
        tk.Label(row1, textvariable=self.template_filename,
                font=("微软雅黑", 9), fg='#666', bg='white').pack(side=tk.LEFT)

        # 内部保存完整路径
        self.template_path = tk.StringVar()

        # 第二行：工作路径及相关按钮
        row2 = tk.Frame(top_section, bg='white')
        row2.pack(fill=tk.X, pady=(0, 0))

        tk.Label(row2, text="选择工作目录", font=("微软雅黑", 11, "bold"),
                bg='white').pack(side=tk.LEFT, padx=(0, 10))

        self.work_path = tk.StringVar()
        work_entry = tk.Entry(row2, textvariable=self.work_path,
                               font=("微软雅黑", 9), relief=tk.SOLID, bd=1, width=35)
        work_entry.pack(side=tk.LEFT, ipady=3)

        RoundedButton(row2, text="选择工作目录", command=self.browse_work_path,
                     bg="#F1EFC0", hover_bg="#E4E1C0", font=("微软雅黑", 9, "bold"),
                     width=100, height=32, corner_radius=10).pack(side=tk.LEFT, padx=(8, 0))

        # 工作目录右侧：填充所有图片、填充所有文本
        RoundedButton(row2, text="填充所有图片", command=self.select_all_images,
                     bg='#E8F4E8', hover_bg='#D4E8D4', font=("微软雅黑", 9, "bold"),
                     width=110, height=32, corner_radius=10).pack(side=tk.LEFT, padx=(10, 0))

        RoundedButton(row2, text="填充所有文本", command=self.fill_all_text,
                     bg='#E8F4FF', hover_bg='#D0E0FF', font=("微软雅黑", 9, "bold"),
                     width=110, height=32, corner_radius=10).pack(side=tk.LEFT, padx=(8, 0))

        # ========== 中间区域：左侧图片列表 + 右侧预览（16:9） ==========
        middle_row = tk.Frame(main_container, bg='white')
        middle_row.pack(fill=tk.BOTH, expand=True, pady=(15, 0))

        # 左侧：图片列表（固定宽度，固定高度）
        list_section = tk.Frame(middle_row, bg='white', width=480, height=372)
        list_section.pack(side=tk.LEFT, fill=tk.BOTH, expand=False)
        list_section.pack_propagate(False)  # 固定大小

        # 标题和按钮在同一行
        list_header = tk.Frame(list_section, bg='white')
        list_header.pack(fill=tk.X, pady=(0, 0))

        tk.Label(list_header, text="图片列表", font=("微软雅黑", 11, "bold"),
                bg='white').pack(side=tk.LEFT)

        # "添加单个图片"按钮
        RoundedButton(list_header, text="+ 添加单个图片", command=self.add_single_image,
                     bg='#E8F4E8', hover_bg='#D4E8D4', font=("微软雅黑", 9),
                     width=130, height=30, corner_radius=10).pack(side=tk.LEFT, padx=(15, 0))

        # "+添加文本"按钮
        RoundedButton(list_header, text="+ 添加文本", command=self.add_text_entry,
                     bg="#F7F3D4", hover_bg="#D8D6B9", font=("微软雅黑", 9),
                     width=130, height=30, corner_radius=10).pack(side=tk.LEFT, padx=(8, 0))

        # "保存当前布局"按钮
        RoundedButton(list_header, text="保存当前布局", command=self.save_current_as_mode,
                     bg='#C8E6C9', hover_bg='#AED6B1', font=("微软雅黑", 9, "bold"),
                     width=120, height=30, corner_radius=10).pack(side=tk.LEFT, padx=(8, 0))

        # 图片列表滚动区域，灰色边框，固定高度362px（减去标题和间距）
        list_canvas_frame = tk.Frame(list_section, bg='white',
                                    highlightbackground='#E0E0E0', highlightthickness=1)
        list_canvas_frame.pack(fill=tk.BOTH, expand=True)

        canvas = tk.Canvas(list_canvas_frame, bg='white', highlightthickness=0)
        # 滚动条（使用加粗样式）
        scrollbar = ttk.Scrollbar(list_canvas_frame, orient="vertical", style='Vertical.TScrollbar', command=canvas.yview)
        self.scrollable_frame = tk.Frame(canvas, bg='white')

        self.scrollable_frame.bind(
            "<Configure>",
            lambda e: canvas.configure(scrollregion=canvas.bbox("all"))
        )

        canvas.create_window((0, 0), window=self.scrollable_frame, anchor="nw")
        canvas.configure(yscrollcommand=scrollbar.set)

        # 鼠标滚轮
        def _on_mousewheel(event):
            canvas.yview_scroll(int(-1*(event.delta/120)), "units")
        canvas.bind_all("<MouseWheel>", _on_mousewheel)

        canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

        # 右侧：布局预览（16:9，固定高度与左侧一致）
        preview_section = tk.Frame(middle_row, bg='white', height=372)
        preview_section.pack(side=tk.RIGHT, padx=(10, 0), fill=tk.BOTH, expand=True)
        preview_section.pack_propagate(False)  # 固定高度

        # 标题行
        preview_header = tk.Frame(preview_section, bg='white')
        preview_header.pack(fill=tk.X, pady=(0, 5))

        tk.Label(preview_header, text="布局预览窗口", font=("微软雅黑", 11, "bold"),
                bg='white').pack(side=tk.LEFT)

        tk.Label(preview_header, text="（16：9）", font=("微软雅黑", 8),
                fg='gray', bg='white').pack(side=tk.LEFT, padx=(8, 0))

        # 16:9 预览画布，带灰色边框，填充剩余空间
        preview_canvas_wrapper = tk.Frame(preview_section, bg='white',
                                         highlightbackground='#E0E0E0', highlightthickness=1)
        preview_canvas_wrapper.pack(fill=tk.BOTH, expand=True)
        self.preview_canvas = LayoutPreviewCanvas(preview_canvas_wrapper, width=667, height=372)
        self.preview_canvas.pack(fill=tk.BOTH, expand=True)

        # ========== 确认插图按钮（左移10px） ==========
        btn_container = tk.Frame(preview_section, bg='white')
        btn_container.pack(side=tk.RIGHT, pady=(5, 0), padx=(10, 0))
        RoundedButton(btn_container, text="确认插图", command=self.generate_ppt,
                     bg='#C8E6C9', hover_bg='#AED6B1', font=("微软雅黑", 13, "bold"),
                     width=160, height=45, corner_radius=10).pack()

        # 预览下方的提示信息区域
        info_label = tk.Label(preview_section, textvariable=self.preview_info_var, fg='green',
                bg='white', font=("微软雅黑", 8), anchor=tk.W,
                wraplength=666, justify=tk.LEFT)
        info_label.pack(fill=tk.X, pady=(8, 0))

        # 初始化
        self.update_preview()

    def browse_template(self):
        """浏览并选择模板PPT"""
        filename = filedialog.askopenfilename(
            title="选择模板PPT",
            filetypes=[("PowerPoint文件", "*.pptx"), ("所有文件", "*.*")]
        )
        if filename:
            self.template_path.set(filename)
            # 只显示文件名
            filename_only = os.path.basename(filename)
            self.template_filename.set(filename_only)
            self.info_hint.set(f"已选择模板: {filename_only}")

            # 如果当前选中的模式有幻灯片索引，应用该模式
            mode = self.current_mode.get()
            if mode != "自定义":
                self.apply_mode_with_template()

    def delete_layout_mode(self):
        """删除布局模式"""
        # 获取可删除的布局列表（排除"自定义"）
        deletable_modes = [m for m in self.preset_modes.keys() if m != "自定义"]
        if not deletable_modes:
            messagebox.showinfo("提示", "没有可删除的布局模式！")
            return

        # 让用户选择要删除的布局
        dialog = tk.Toplevel(self.root)
        dialog.title("删除布局模式")
        dialog.geometry("350x300")
        dialog.configure(bg='white')
        dialog.transient(self.root)
        dialog.grab_set()

        # 居中显示
        x = self.root.winfo_x() + (self.root.winfo_width() - 350) // 2
        y = self.root.winfo_y() + (self.root.winfo_height() - 300) // 2
        dialog.geometry(f"+{x}+{y}")

        content = tk.Frame(dialog, bg='white', padx=20, pady=20)
        content.pack(fill=tk.BOTH, expand=True)

        tk.Label(content, text="选择要删除的布局模式：", font=("微软雅黑", 11, "bold"),
                 bg='white').pack(pady=(0, 10))

        # 创建列表框
        mode_list = tk.Listbox(content, height=8,
                             font=("微软雅黑", 10), bg='white', relief=tk.SOLID, bd=1)
        for mode in deletable_modes:
            mode_list.insert(tk.END, mode)
        mode_list.pack(fill=tk.X, pady=(0, 15))

        btn_row = tk.Frame(content, bg='white')
        btn_row.pack(fill=tk.X)

        def confirm_delete():
            selection = mode_list.curselection()
            if not selection:
                messagebox.showwarning("提示", "请选择要删除的布局模式！")
                return

            mode_to_delete = deletable_modes[selection[0]]

            if messagebox.askyesno("确认删除", f"确定要删除布局模式'{mode_to_delete}'吗？"):
                del self.preset_modes[mode_to_delete]
                save_custom_modes(self.preset_modes)

                self.mode_combo['values'] = list(self.preset_modes.keys())
                if self.current_mode.get() == mode_to_delete:
                    self.current_mode.set("自定义")

                self.update_preview()
                self.info_hint.set(f"已删除布局模式: {mode_to_delete}")
                dialog.destroy()

        RoundedButton(btn_row, text="确定删除", command=confirm_delete,
                     bg='#FFE0E0', hover_bg='#FFD0D0', font=("微软雅黑", 10),
                     width=120, height=32, corner_radius=10).pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 10))

        RoundedButton(btn_row, text="取消", command=dialog.destroy,
                     bg='#F0F0F0', hover_bg='#E0E0E0', font=("微软雅黑", 10),
                     width=120, height=32, corner_radius=10).pack(side=tk.LEFT, fill=tk.X, expand=True)

    def browse_work_path(self):
        """选择工作路径"""
        dirname = filedialog.askdirectory(title="选择工作目录")
        if dirname:
            self.work_path.set(dirname)
            self.info_hint.set(f"工作路径: {dirname}")

    def apply_mode_with_template(self):
        """应用布局模式（包含模板信息）"""
        mode_name = self.current_mode.get()
        if mode_name == "自定义":
            return

        mode_config = self.preset_modes.get(mode_name)
        if not mode_config:
            return

        # 如果模式包含模板信息，自动设置模板路径
        if mode_config.get("template_file"):
            template_file = mode_config["template_file"]
            if os.path.exists(template_file):
                self.template_path.set(template_file)
                # 更新模板文件名显示
                self.template_filename.set(os.path.basename(template_file))
            else:
                # 尝试使用默认模板文件
                prog_dir = os.path.dirname(os.path.abspath(__file__))
                default_template = os.path.join(prog_dir, DEFAULT_TEMPLATE_FILE)
                if os.path.exists(default_template):
                    self.template_path.set(default_template)
                    self.template_filename.set(os.path.basename(default_template))

    def on_mode_change(self, event=None):
        """当布局模式改变时更新预览、设置模板并应用模式"""
        self.update_preview()  # 更新预览
        self.apply_mode_with_template()  # 设置模板
        self.apply_mode()  # 应用图片布局（填充图片列表）

    def update_preview(self):
        """更新布局预览"""
        mode = self.current_mode.get()
        layouts = self.preset_modes[mode]["layouts"]
        text_layouts = self.preset_modes[mode].get("text_layouts", [])
        self.preview_canvas.draw_layout(layouts, text_layouts)

    def apply_mode(self):
        """应用选中的布局模式（包含模板信息和图片位置）"""
        mode_name = self.current_mode.get()
        mode_config = self.preset_modes[mode_name]
        layouts = mode_config["layouts"]
        text_layouts = mode_config.get("text_layouts", [])

        if not layouts and not text_layouts:
            return

        # 自动设置模板路径（如果模式中有模板信息）
        self.apply_mode_with_template()

        # 清空现有条目
        self.clear_all_entries()

        # 根据模式创建图片条目
        for i, layout in enumerate(layouts):
            entry = ImageEntry(self.scrollable_frame, self.remove_image_entry, i, app_master=self)
            entry.set_layout(layout)
            self.image_entries.append(entry)

        # 根据模式创建文本条目
        for i, text_layout in enumerate(text_layouts):
            entry = TextEntry(self.scrollable_frame, self.remove_text_entry, i, app_master=self)
            entry.set_layout(text_layout)
            self.text_entries.append(entry)

        slide_info = ""
        if mode_config.get("slide_index") is not None:
            slide_num = mode_config["slide_index"] + 1
            slide_info = f"，使用模板第{slide_num}页"

        # 显示提示信息：已选中的模式名和图片数量
        text_info = f"，{len(text_layouts)}个文本" if text_layouts else ""
        self.preview_info_var.set(f"已选中布局模式'{mode_name}'，{len(layouts)}张图{text_info}待插入")

    def select_all_images(self):
        """全选工作路径下所有可用图片"""
        work_dir = self.work_path.get()
        if not work_dir:
            messagebox.showwarning("提示", "请先选择工作路径！")
            return

        if not os.path.exists(work_dir):
            messagebox.showwarning("提示", f"工作路径不存在: {work_dir}")
            return

        # 获取所有图片文件
        image_extensions = ['.jpg', '.jpeg', '.png', '.bmp', '.gif']
        image_files = []

        try:
            for filename in os.listdir(work_dir):
                if os.path.splitext(filename.lower())[1] in image_extensions:
                    image_files.append(filename)
        except Exception as e:
            messagebox.showerror("错误", f"读取文件失败: {str(e)}")
            return

        if not image_files:
            messagebox.showinfo("提示", f"工作路径 {work_dir} 下没有找到图片文件")
            return

        # 按文件名排序
        image_files.sort()

        # 提示用户
        if not self.image_entries:
            result = messagebox.askyesno("确认",
                f"找到 {len(image_files)} 张图片\n\n是否创建图片列表？")
            if not result:
                return
        else:
            expected_count = len(self.image_entries)
            if len(image_files) > expected_count:
                result = messagebox.askyesno("确认",
                    f"找到 {len(image_files)} 张图片\n当前布局需要 {expected_count} 张\n\n是否只使用前 {expected_count} 张？")
                if not result:
                    return
                image_files = image_files[:expected_count]
            elif len(image_files) < expected_count:
                result = messagebox.askyesno("确认",
                    f"找到 {len(image_files)} 张图片\n当前布局需要 {expected_count} 张\n\n是否继续？（缺少的位置将保持空白）")
                if not result:
                    return

        # 清空现有条目（如果需要）
        if not self.image_entries:
            self.clear_all_entries()

        # 创建或填入图片条目
        for i, filename in enumerate(image_files):
            if i < len(self.image_entries):
                # 已有序目，只填文件名
                self.image_entries[i].image_filename.set(filename)
            else:
                # 创建新条目
                entry = ImageEntry(self.scrollable_frame, self.remove_image_entry, i, app_master=self)
                entry.image_filename.set(filename)
                self.image_entries.append(entry)

        self.list_info_var.set(f"已加载 {len(image_files)} 张图片，当前共 {len(self.image_entries)} 个")
        self.preview_info_var.set(f"已填充 {len(image_files)} 张图片")

    def add_single_image(self):
        """添加单个图片（选择文件名）"""
        work_dir = self.work_path.get()
        if not work_dir:
            messagebox.showwarning("提示", "请先选择工作路径！")
            return

        if not os.path.exists(work_dir):
            messagebox.showwarning("提示", f"工作路径不存在: {work_dir}")
            return

        # 在工作路径下选择图片
        filename = filedialog.askopenfilename(
            title="选择图片",
            initialdir=work_dir,
            filetypes=[
                ("图片文件", "*.jpg *.jpeg *.png *.bmp *.gif"),
                ("所有文件", "*.*")
            ]
        )

        if filename:
            index = len(self.image_entries)
            entry = ImageEntry(self.scrollable_frame, self.remove_image_entry, index, app_master=self)
            entry.image_filename.set(os.path.basename(filename))
            self.image_entries.append(entry)
            self.list_info_var.set(f"已添加图片，当前共 {len(self.image_entries)} 个")

    def add_text_entry(self):
        """添加一个文本条目"""
        index = len(self.text_entries)
        entry = TextEntry(self.scrollable_frame, self.remove_text_entry, index, app_master=self)
        self.text_entries.append(entry)
        self.list_info_var.set(f"已添加文本，当前共 {len(self.text_entries)} 个")

    def save_current_as_mode(self):
        if not self.image_entries and not self.text_entries:
            self.list_info_var.set("请先添加图片或文本并配置位置！")
            return

        # 简化的对话框
        dialog = tk.Toplevel(self.root)
        dialog.title("保存为布局模式")
        dialog.geometry("500x200")  # 增大对话框尺寸
        dialog.transient(self.root)
        dialog.grab_set()
        dialog.configure(bg='white')

        # 将对话框居中显示在主窗口
        dialog.update_idletasks()  # 确保获取窗口尺寸
        x = self.root.winfo_x() + (self.root.winfo_width() - 500) // 2
        y = self.root.winfo_y() + (self.root.winfo_height() - 200) // 2
        dialog.geometry(f"+{x}+{y}")

        content_frame = tk.Frame(dialog, bg='white', padx=30, pady=30)
        content_frame.pack(fill=tk.BOTH, expand=True)

        tk.Label(content_frame, text="模式名称:", bg='white', font=("微软雅黑", 11)).pack(anchor=tk.W, pady=(0, 8))
        name_var = tk.StringVar()
        tk.Entry(content_frame, textvariable=name_var, width=40, font=("微软雅黑", 10)).pack(pady=(0, 20))

        def save_mode():
            name = name_var.get().strip()

            if not name:
                self.list_info_var.set("请输入模式名称！")
                return

            if name == "自定义":
                self.list_info_var.set("不能使用'自定义'作为模式名称！")
                return

            # 收集当前图片布局
            layouts = []
            for entry in self.image_entries:
                try:
                    layout = {
                        "slide": 0,  # 默认
                        "left": float(entry.left_var.get()),
                        "top": float(entry.top_var.get()),
                    }
                    if entry.width_var.get().strip():
                        layout["width"] = float(entry.width_var.get())
                    if entry.height_var.get().strip():
                        layout["height"] = float(entry.height_var.get())

                    layouts.append(layout)
                except Exception as e:
                    self.list_info_var.set(f"读取图片配置失败: {str(e)}")
                    return

            # 收集当前文本布局
            text_layouts = []
            for entry in self.text_entries:
                try:
                    text_layout = entry.get_config()
                    text_layouts.append(text_layout)
                except Exception as e:
                    self.list_info_var.set(f"读取文本配置失败: {str(e)}")
                    return

            if not layouts and not text_layouts:
                self.list_info_var.set("没有有效的配置可保存！")
                return

            # 保存模式（包含模板信息和文本布局）
            mode_data = {
                "description": f"{len(layouts)}张图片布局",
                "template_file": self.template_path.get() if self.template_path.get() else None,
                "slide_index": 0,  # 默认使用第1页
                "layouts": layouts,
                "text_layouts": text_layouts  # 新增：文本布局
            }

            if text_layouts:
                mode_data["description"] += f"，{len(text_layouts)}个文本"

            # 如果设置了模板文件，询问使用哪一页
            if self.template_path.get() and os.path.exists(self.template_path.get()):
                try:
                    prs = Presentation(self.template_path.get())
                    total_slides = len(prs.slides)
                    slide_num = simpledialog.askinteger(
                        "选择模板页",
                        f"模板文件共有 {total_slides} 页幻灯片\n\n请输入使用第几页作为模板（1-{total_slides}）：",
                        parent=dialog,
                        minvalue=1,
                        maxvalue=total_slides,
                        initialvalue=1
                    )
                    if slide_num:
                        mode_data["slide_index"] = slide_num - 1
                        mode_data["description"] = f"{len(layouts)}张图片布局（模板第{slide_num}页）"
                        if text_layouts:
                            mode_data["description"] += f"，{len(text_layouts)}个文本"
                except:
                    pass

            self.preset_modes[name] = mode_data

            if save_custom_modes(self.preset_modes):
                # 更新下拉列表
                self.mode_combo['values'] = list(self.preset_modes.keys())
                self.current_mode.set(name)
                self.on_mode_change()
                self.list_info_var.set(f"已保存模式'{name}'，包含 {len(layouts)} 个图片位置，{len(text_layouts)} 个文本")
                dialog.destroy()

        btn_frame = tk.Frame(content_frame, bg='white')
        btn_frame.pack()

        RoundedButton(btn_frame, text="保存", command=save_mode,
                     bg='#C8E6C9', hover_bg='#AED6B1', font=("微软雅黑", 10, "bold"),
                     width=120, height=35, corner_radius=10).pack()

    def remove_image_entry(self, entry):
        """移除图片条目"""
        if entry in self.image_entries:
            self.image_entries.remove(entry)
            self.list_info_var.set(f"已删除图片条目，当前共 {len(self.image_entries)} 个")

    def remove_text_entry(self, entry):
        """移除文本条目"""
        if entry in self.text_entries:
            self.text_entries.remove(entry)
            self.list_info_var.set(f"已删除文本条目，当前共 {len(self.text_entries)} 个")

    def clear_all_entries(self):
        """清空所有条目"""
        for entry in self.image_entries[:]:
            entry.frame.destroy()
        self.image_entries.clear()
        for entry in self.text_entries[:]:
            entry.frame.destroy()
        self.text_entries.clear()
        self.list_info_var.set("已清空所有条目")

    def generate_ppt(self):
        """生成PPT文件"""
        if not self.image_entries and not self.text_entries:
            self.preview_info_var.set("请至少添加一个图片或文本！")
            return

        template = self.template_path.get()
        work_dir = self.work_path.get()

        if not template:
            self.preview_info_var.set("请选择模板PPT！")
            return

        if not work_dir:
            self.preview_info_var.set("请选择工作路径！")
            return

        # 自动生成输出文件名：布局模式+时-分-秒
        mode_name = self.current_mode.get()
        current_time = datetime.now().strftime("%H-%M-%S")
        output_filename = f"{mode_name}_{current_time}.pptx"

        # 构建输出文件路径
        output = os.path.join(work_dir, output_filename)

        try:
            self.preview_info_var.set("正在插入图片...")
            self.root.update()

            # 打开模板PPT
            if os.path.exists(template):
                prs = Presentation(template)

                # 获取当前布局模式的幻灯片索引（如果有的话）
                mode_name = self.current_mode.get()
                mode_config = self.preset_modes.get(mode_name, {})
                source_slide_index = mode_config.get("slide_index", 0)

                # 检查幻灯片索引是否有效
                if source_slide_index < 0 or source_slide_index >= len(prs.slides):
                    source_slide_index = 0

                mode_text = f"基于模板: {os.path.basename(template)}，第{source_slide_index + 1}页"
            else:
                self.preview_info_var.set(f"模板文件不存在: {template}")
                return

            # 确保源幻灯片存在
            while len(prs.slides) <= source_slide_index:
                prs.slides.add_slide(prs.slide_layouts[0])

            # 复制源幻灯片创建新幻灯片
            source_slide = prs.slides[source_slide_index]
            new_slide = prs.slides.add_slide(source_slide.slide_layout)

            # 只复制非图片元素（跳过图片，避免重复）
            for shape in source_slide.shapes:
                try:
                    # 跳过图片类型的形状
                    if shape.shape_type == 13:  # 13 = MSO_SHAPE_TYPE.PICTURE
                        continue

                    el = shape.element
                    newel = deepcopy(el)
                    new_slide.shapes._spTree.append(newel)
                except:
                    pass

            # 删除所有旧幻灯片，只保留新创建的幻灯片
            slides_to_delete = list(range(len(prs.slides) - 1))
            for idx in reversed(slides_to_delete):
                rId = prs.slides._sldIdLst[idx].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[idx]

            # 处理每个图片
            success_count = 0
            errors = []

            for i, entry in enumerate(self.image_entries):
                try:
                    config = entry.get_config()

                    # 检查图片文件名
                    if not config['filename']:
                        errors.append(f"图片{i+1}: 未选择图片文件名")
                        continue

                    # 构建完整路径（工作路径 + 文件名）
                    image_path = os.path.join(work_dir, config['filename'])

                    if not os.path.exists(image_path):
                        errors.append(f"图片{i+1}: 找不到文件 {config['filename']}")
                        continue

                    # 转换单位
                    left = Cm(config['left'])
                    top = Cm(config['top'])
                    width = Cm(config['width']) if 'width' in config else None
                    height = Cm(config['height']) if 'height' in config else None

                    # 插入图片到幻灯片
                    if width and height:
                        new_slide.shapes.add_picture(image_path, left, top, width=width, height=height)
                    elif width:
                        new_slide.shapes.add_picture(image_path, left, top, width=width)
                    elif height:
                        new_slide.shapes.add_picture(image_path, left, top, height=height)
                    else:
                        new_slide.shapes.add_picture(image_path, left, top)

                    success_count += 1

                except Exception as e:
                    errors.append(f"图片{i+1}: {str(e)}")
                    continue

            # 处理每个文本（根据关键词从不同文件中提取）
            text_success_count = 0
            text_errors = []

            for i, entry in enumerate(self.text_entries):
                try:
                    config = entry.get_config()
                    keyword = config.get('keyword', '').strip()

                    # 根据关键词搜索文件
                    matched_file = None
                    if keyword:
                        # 搜索包含关键词的文件
                        try:
                            for filename in os.listdir(work_dir):
                                if keyword in filename:
                                    filepath = os.path.join(work_dir, filename)
                                    if os.path.isfile(filepath):
                                        matched_file = filename
                                        break
                        except Exception as e:
                            text_errors.append(f"文本{i+1}: 搜索文件失败 - {str(e)}")
                            continue
                    else:
                        # 如果没有关键词，搜索所有文本文件
                        text_files = []
                        text_extensions = ['.txt', '.csv', '.log', '.dat', '.json', '.xml']
                        try:
                            for filename in os.listdir(work_dir):
                                # 检查标准扩展名
                                if os.path.splitext(filename.lower())[1] in text_extensions:
                                    text_files.append(filename)
                                # 检查.o数字格式（如 .o2343908）
                                elif re.search(r'\.o\d+$', filename):
                                    text_files.append(filename)
                        except Exception as e:
                            text_files = []

                        # 如果没有找到文本文件，也尝试读取所有文件，判断是否为文本
                        if not text_files:
                            try:
                                for filename in os.listdir(work_dir):
                                    filepath = os.path.join(work_dir, filename)
                                    if os.path.isfile(filepath):
                                        # 先检查是否是.o数字格式
                                        if not re.search(r'\.o\d+$', filename):
                                            # 尝试读取文件前100字节判断是否为文本
                                            try:
                                                with open(filepath, 'rb') as f:
                                                    chunk = f.read(100)
                                                    # 检查是否包含大量控制字符（非文本）
                                                    text_chars = 0
                                                    for byte in chunk:
                                                        if byte >= 32 or byte in [9, 10, 13]:  # 空格及可打印字符、制表符、换行符
                                                            text_chars += 1
                                                    if text_chars / max(len(chunk), 1) > 0.7:  # 70%以上是文本字符
                                                        text_files.append(filename)
                                            except:
                                                continue
                            except Exception as e:
                                pass

                        # 按文件名排序，使用第一个文件
                        text_files.sort()
                        if text_files:
                            matched_file = text_files[0]

                    if not matched_file:
                        if keyword:
                            text_errors.append(f"文本{i+1}: 找不到包含关键词'{keyword}'的文件")
                        else:
                            text_errors.append(f"文本{i+1}: 找不到文本文件")
                        continue

                    text_path = os.path.join(work_dir, matched_file)

                    if not os.path.exists(text_path):
                        text_errors.append(f"文本{i+1}: 找不到文件 {matched_file}")
                        continue

                    # 尝试多种编码读取文本文件
                    line_content = None
                    encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16', 'latin-1']
                    for encoding in encodings:
                        try:
                            with open(text_path, 'r', encoding=encoding) as f:
                                lines = f.readlines()

                            # 获取指定行（行号从1开始）
                            line_index = config['line_number'] - 1
                            if 0 <= line_index < len(lines):
                                line_content = lines[line_index].strip()
                                break
                            else:
                                # 编码读取成功但行数不够，继续尝试其他编码
                                continue
                        except UnicodeDecodeError:
                            continue
                        except Exception:
                            continue

                    if line_content is None:
                        # 所有编码都尝试失败，使用二进制方式读取并尝试解码
                        try:
                            with open(text_path, 'rb') as f:
                                lines = f.read().decode('utf-8', errors='ignore').split('\n')
                            line_index = config['line_number'] - 1
                            if 0 <= line_index < len(lines):
                                line_content = lines[line_index].strip()
                            else:
                                text_errors.append(f"文本{i+1}: 第{config['line_number']}行不存在")
                                continue
                        except Exception as e:
                            text_errors.append(f"文本{i+1}: 读取文件失败 - {str(e)}")
                            continue

                    # 解析需要读取的列（多个列用逗号分隔）
                    file_cols_str = config['file_cols']
                    col_values = []
                    try:
                        col_numbers = [int(x.strip()) for x in file_cols_str.split(',')]
                        for col_num in col_numbers:
                            col_idx = col_num - 1  # 转换为索引
                            if 0 <= col_idx < len(line_content.split()):
                                # 按空格或制表符分割行内容
                                words = line_content.split()
                                if col_idx < len(words):
                                    col_values.append(words[col_idx])
                    except Exception as e:
                        text_errors.append(f"文本{i+1}: 解析列号失败 - {str(e)}")
                        continue

                    # 如果没有成功获取任何列的值，整行作为默认值
                    if not col_values:
                        text_content = line_content
                        # 格式化单列数据
                        text_content = format_text(text_content)
                    else:
                        # 使用/分隔多个列的值，每列都格式化
                        formatted_values = [format_text(val) for val in col_values]
                        text_content = '/'.join(formatted_values)

                    # 添加文本框（通过左、上坐标定位）
                    left = Cm(config['left'])
                    top = Cm(config['top'])

                    # 添加文本框
                    text_box = new_slide.shapes.add_textbox(left, top, width=Cm(5), height=Cm(1))

                    # 设置文本内容
                    text_frame = text_box.text_frame
                    text_frame.word_wrap = False

                    for paragraph in text_frame.paragraphs:
                        paragraph.text = text_content
                        # 文本默认左对齐
                        paragraph.alignment = PP_ALIGN.LEFT

                        for run in paragraph.runs:
                            run.font.name = 'LiciumFont 2022'
                            run.font.size = Pt(20)
                            run.font.color.rgb = RGBColor(0, 0, 0)
                            # 字体不加粗
                            run.font.bold = False

                    text_success_count += 1

                except Exception as e:
                    text_errors.append(f"文本{i+1}: {str(e)}")
                    continue

            # 保存到输出路径
            prs.save(output)

            # 显示结果
            if success_count == 0 and text_success_count == 0:
                # 如果没有成功插入任何内容，直接显示警告
                result_msg = f"警告: 没有成功插入任何内容\n"
                all_errors = errors + text_errors
                if all_errors:
                    result_msg += f"{len(all_errors)} 个错误: " + "; ".join(all_errors[:3])
                    if len(all_errors) > 3:
                        result_msg += f"... 还有 {len(all_errors)-3} 个"
            else:
                # 有内容成功插入，显示成功信息
                text_info = f"，{text_success_count}个文本" if text_success_count > 0 else ""
                result_msg = f"（成功插入）{mode_text}\n成功插入 {success_count}/{len(self.image_entries)} 张图片{text_info}\n保存位置: {output}"

                all_errors = errors + text_errors
                if all_errors:
                    result_msg += f"\n警告: {len(all_errors)} 个错误: " + "; ".join(all_errors[:3])
                    if len(all_errors) > 3:
                        result_msg += f"... 还有 {len(all_errors)-3} 个"

            self.preview_info_var.set(result_msg)

        except Exception as e:
            self.preview_info_var.set(f"操作失败: {str(e)}")

    def fill_all_text(self):
        """填充所有文本到模板PPT（不生成新PPT，直接在模板中填充文本）"""
        if not self.text_entries:
            messagebox.showinfo("提示", "没有文本需要填充")
            return

        template = self.template_path.get()
        work_dir = self.work_path.get()

        if not template:
            messagebox.showwarning("提示", "请先选择模板PPT！")
            return

        if not work_dir:
            messagebox.showwarning("提示", "请先选择工作路径！")
            return

        try:
            self.preview_info_var.set("正在填充文本...")
            self.root.update()

            # 打开模板PPT
            if os.path.exists(template):
                prs = Presentation(template)

                # 获取当前布局模式的幻灯片索引
                mode_name = self.current_mode.get()
                mode_config = self.preset_modes.get(mode_name, {})
                slide_index = mode_config.get("slide_index", 0)

                # 检查幻灯片索引是否有效
                if slide_index < 0 or slide_index >= len(prs.slides):
                    messagebox.showwarning("提示", f"幻灯片索引超出范围（共{len(prs.slides)}页）")
                    return

                slide = prs.slides[slide_index]
            else:
                messagebox.showerror("错误", f"模板文件不存在: {template}")
                return

            success_count = 0
            errors = []

            for i, entry in enumerate(self.text_entries):
                try:
                    config = entry.get_config()
                    keyword = config.get('keyword', '').strip()

                    # 根据关键词搜索文件
                    matched_file = None
                    if keyword:
                        # 搜索包含关键词的文件
                        try:
                            for filename in os.listdir(work_dir):
                                if keyword in filename:
                                    filepath = os.path.join(work_dir, filename)
                                    if os.path.isfile(filepath):
                                        matched_file = filename
                                        break
                        except Exception as e:
                            errors.append(f"文本{i+1}: 搜索文件失败 - {str(e)}")
                            continue
                    else:
                        # 如果没有关键词，搜索所有文本文件
                        text_files = []
                        text_extensions = ['.txt', '.csv', '.log', '.dat', '.json', '.xml']
                        try:
                            for filename in os.listdir(work_dir):
                                # 检查标准扩展名
                                if os.path.splitext(filename.lower())[1] in text_extensions:
                                    text_files.append(filename)
                                # 检查.o数字格式（如 .o2343908）
                                elif re.search(r'\.o\d+$', filename):
                                    text_files.append(filename)
                        except Exception as e:
                            text_files = []

                        # 如果没有找到文本文件，也尝试读取所有文件，判断是否为文本
                        if not text_files:
                            try:
                                for filename in os.listdir(work_dir):
                                    filepath = os.path.join(work_dir, filename)
                                    if os.path.isfile(filepath):
                                        # 先检查是否是.o数字格式
                                        if not re.search(r'\.o\d+$', filename):
                                            # 尝试读取文件前100字节判断是否为文本
                                            try:
                                                with open(filepath, 'rb') as f:
                                                    chunk = f.read(100)
                                                    # 检查是否包含大量控制字符（非文本）
                                                    text_chars = 0
                                                    for byte in chunk:
                                                        if byte >= 32 or byte in [9, 10, 13]:  # 空格及可打印字符、制表符、换行符
                                                            text_chars += 1
                                                    if text_chars / max(len(chunk), 1) > 0.7:  # 70%以上是文本字符
                                                        text_files.append(filename)
                                            except:
                                                continue
                            except Exception as e:
                                pass

                        # 按文件名排序，使用第一个文件
                        text_files.sort()
                        if text_files:
                            matched_file = text_files[0]

                    if not matched_file:
                        if keyword:
                            errors.append(f"文本{i+1}: 找不到包含关键词'{keyword}'的文件")
                        else:
                            errors.append(f"文本{i+1}: 找不到文本文件")
                        continue

                    text_path = os.path.join(work_dir, matched_file)

                    if not os.path.exists(text_path):
                        errors.append(f"文本{i+1}: 找不到文件 {matched_file}")
                        continue

                    # 尝试多种编码读取文本文件
                    line_content = None
                    encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16', 'latin-1']
                    for encoding in encodings:
                        try:
                            with open(text_path, 'r', encoding=encoding) as f:
                                lines = f.readlines()

                            # 获取指定行（行号从1开始）
                            line_index = config['line_number'] - 1
                            if 0 <= line_index < len(lines):
                                line_content = lines[line_index].strip()
                                break
                            else:
                                # 编码读取成功但行数不够，继续尝试其他编码
                                continue
                        except UnicodeDecodeError:
                            continue
                        except Exception:
                            continue

                    if line_content is None:
                        # 所有编码都尝试失败，使用二进制方式读取并尝试解码
                        try:
                            with open(text_path, 'rb') as f:
                                lines = f.read().decode('utf-8', errors='ignore').split('\n')
                            line_index = config['line_number'] - 1
                            if 0 <= line_index < len(lines):
                                line_content = lines[line_index].strip()
                            else:
                                errors.append(f"文本{i+1}: 第{config['line_number']}行不存在")
                                continue
                        except Exception as e:
                            errors.append(f"文本{i+1}: 读取文件失败 - {str(e)}")
                            continue

                    # 解析需要读取的列（多个列用逗号分隔）
                    file_cols_str = config['file_cols']
                    col_values = []
                    try:
                        col_numbers = [int(x.strip()) for x in file_cols_str.split(',')]
                        for col_num in col_numbers:
                            col_idx = col_num - 1  # 转换为索引
                            if 0 <= col_idx < len(line_content.split()):
                                # 按空格或制表符分割行内容
                                words = line_content.split()
                                if col_idx < len(words):
                                    col_values.append(words[col_idx])
                    except Exception as e:
                        errors.append(f"文本{i+1}: 解析列号失败 - {str(e)}")
                        continue

                    # 如果没有成功获取任何列的值，整行作为默认值
                    if not col_values:
                        text_content = line_content
                        # 格式化单列数据
                        text_content = format_text(text_content)
                    else:
                        # 使用/分隔多个列的值，每列都格式化
                        formatted_values = [format_text(val) for val in col_values]
                        text_content = '/'.join(formatted_values)

                    # 添加文本框（通过左、上坐标定位）
                    left = Cm(config['left'])
                    top = Cm(config['top'])

                    # 添加文本框
                    text_box = slide.shapes.add_textbox(left, top, width=Cm(5), height=Cm(1))

                    # 设置文本内容
                    text_frame = text_box.text_frame
                    text_frame.word_wrap = False

                    for paragraph in text_frame.paragraphs:
                        paragraph.text = text_content
                        # 文本默认左对齐
                        paragraph.alignment = PP_ALIGN.LEFT

                        for run in paragraph.runs:
                            run.font.name = 'LiciumFont 2022'
                            run.font.size = Pt(20)
                            run.font.color.rgb = RGBColor(0, 0, 0)
                            # 字体不加粗
                            run.font.bold = False

                    success_count += 1

                except Exception as e:
                    errors.append(f"文本{i+1}: {str(e)}")
                    continue

            # 不保存到输出路径（只需在模板中填充文本即可）
            # prs.save(output)

            # 显示结果
            result_msg = f"成功填充 {success_count}/{len(self.text_entries)} 个文本"
            if errors:
                result_msg += f"\n警告: {len(errors)} 个错误: " + "; ".join(errors[:3])
                if len(errors) > 3:
                    result_msg += f"... 还有 {len(errors)-3} 个"

            self.preview_info_var.set(result_msg)

        except Exception as e:
            self.preview_info_var.set(f"操作失败: {str(e)}")
            messagebox.showerror("错误", f"操作失败: {str(e)}")


def main():
    """主函数"""
    root = tk.Tk()

    # 设置样式
    style = ttk.Style()
    style.theme_use('clam')

    app = PPTImageInserterGUI(root)
    root.mainloop()


if __name__ == "__main__":
    main()
